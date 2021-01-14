import sys
import glob
import os
import cv2
import pandas as pd
import datetime
from pdf2image import convert_from_path
from os.path import isfile, join
from os import listdir
from PyQt5.QtWidgets import *
from PyQt5.QtCore import QCoreApplication, Qt
from PyQt5.QtGui import *
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from openpyxl import Workbook

try: 
    from PIL import Image 
except ImportError: 
    import Image 
import pytesseract 

path = r'./source/'
img_dir = r'./result/'
obj_dir = r'./object/'
find_dir = r'./find/'

methods = ['cv2.TM_CCOEFF_NORMED']#, 'cv2.TM_CCORR_NORMED', 'cv2.TM_SQDIFF_NORMED'] # 찾는 방법
dpi_list = ['100','200','300','400','800','1600']

custom_config = r'-l eng+kor'

class ImageConverter:
    def convertPDFtoPNG(self,path):
        self.path = path
        self.pdf_dir = glob.glob(self.path + "/*")
        for pdf_ in self.pdf_dir:
            pages = convert_from_path(pdf_, dpi=800)
            for i, page in enumerate(pages):
                page.save(f'{img_dir+pdf_[len(path):-4]}_page{i+1:0>2d}.png','PNG')
                print(f'{pdf_[len(path):-4]}_page{i+1:0>2d}.png saved...')
            print('Done !')

class MyCheckBox(QCheckBox): 
    def __init__(self, item): 
        super().__init__() 
        self.item = item 
        self.mycheckvalue = 0 # 0 --> unchecked, 2 --> checked 
        self.stateChanged.connect(self.__checkbox_change) 
        self.stateChanged.connect(self.item.my_setdata) # checked 여부로 정렬을 하기위한 data 저장 
    
    def __checkbox_change(self, checkvalue): 
        self.mycheckvalue = checkvalue 
        print("checkbox row= ", self.get_row()) 
    
    def get_row(self): 
        return self.item.row()

class MyQTableWidgetItemCheckBox(QTableWidgetItem): 
    def __init__(self): 
        super().__init__() 
        self.setData(Qt.UserRole, 0) 
        
    def __lt__(self, other): 
        return self.data(Qt.UserRole) < other.data(Qt.UserRole) 
        
    def my_setdata(self, value): 
        self.setData(Qt.UserRole, value)

class MyWindow(QWidget):  
    def __init__(self):
        super().__init__()
        self.setupUI()

    def setupUI(self):
        layout = QGridLayout()

        # Drawing group
        self.drawing_group = QGroupBox('drawing')
        self.drawing_layout = QGridLayout()
        self.drawing_loc_le = QLineEdit()
        self.drawing_loc_le.setPlaceholderText('./source/')
        self.drawing_combo = QComboBox(self)
        self.drawing_combo.addItems(dpi_list)
        self.drawing_combo.setCurrentIndex(1)
        self.drawing_button_select = QPushButton("파일 선택")
        self.drawing_button_select.clicked.connect(self.selectButtonClicked)
        self.drawing_button_change_image = QPushButton("이미지 변환")
        self.drawing_button_change_image.clicked.connect(self.changeToImageButtonClicked)
        self.drawing_tabs = QTabWidget()
        self.drawing_label = [QLabel('PDF-image',self) for i in range(10)]
        for idx, label in enumerate(self.drawing_label):
            label.setFixedSize(600,400)
            label.setAlignment(Qt.AlignTop)
            self.drawing_tabs.addTab(label, str(idx+1) + ' ')
        self.drawing_layout.addWidget(self.drawing_loc_le,0,0,1,4)
        self.drawing_layout.addWidget(self.drawing_combo,0,4,1,1)
        self.drawing_layout.addWidget(self.drawing_button_select,0,5,1,1)
        self.drawing_layout.addWidget(self.drawing_button_change_image,0,6,1,1)
        self.drawing_layout.addWidget(self.drawing_tabs,1,0,7,7)
        self.drawing_group.setLayout(self.drawing_layout)
        layout.addWidget(self.drawing_group,0,0,5,1)
        
        # finding & ocr 
        self.finding_group = QGroupBox('find & OCR')
        self.finding_layout = QGridLayout()
        self.finding_le = QLineEdit()
        self.finding_le.setPlaceholderText('./object/')
        self.finding_obj_btn = QPushButton("파일 선택")
        self.finding_obj_btn.clicked.connect(self.findingButtonClicked)
        self.finding_anal_btn = QPushButton("찾아내기")
        self.finding_anal_btn.clicked.connect(self.analyzeButtonClicked)
        self.finding_obj_label = QLabel('object-image',self)
        self.finding_obj_label.setFixedSize(600,350)
        self.finding_obj_label.setAlignment(Qt.AlignTop)
        self.finding_match_label = QLabel('match-image',self)
        self.finding_match_label.setFixedSize(600,350)
        self.finding_match_label.setAlignment(Qt.AlignTop)
        
        self.finding_tabs = QTabWidget()
        self.finding_tabs.addTab(self.finding_obj_label, 'finding_obj')
        self.finding_tabs.addTab(self.finding_match_label, 'found_match')

        self.finding_layout.addWidget(self.finding_le,0,0,1,4)
        self.finding_layout.addWidget(self.finding_obj_btn,0,4,1,1)
        self.finding_layout.addWidget(self.finding_anal_btn,0,5,1,1)
#        self.finding_layout.addWidget(self.finding_obj_label,1,0,6,6)
#        self.finding_layout.addWidget(self.finding_match_label,1,6,6,6)
        self.finding_layout.addWidget(self.finding_tabs,1,0,6,6)
        self.finding_group.setLayout(self.finding_layout)
        layout.addWidget(self.finding_group,5,0,7,1)

        # 중간과정 문자표시 테이블
        self.terminal_group = QGroupBox('teminal')
        self.terminal_layout = QVBoxLayout()
        self.terminal_browser = QTextBrowser()
        self.terminal_browser.setAcceptRichText(True)
        self.terminal_browser.setOpenExternalLinks(True)
        self.terminal_browser.setFixedWidth(500)
        self.terminal_layout.addWidget(self.terminal_browser)
        self.terminal_group.setLayout(self.terminal_layout)
        layout.addWidget(self.terminal_group, 4,1,7,1)

        # function button
        self.func_button_group = QGroupBox('function button')
        self.func_button_layout = QGridLayout()
        self.func_button_clear = QPushButton("Clear")
        self.func_button_clear.clicked.connect(self.clear_text)
        self.func_button_quit = QPushButton("Quit")
        self.func_button_quit.clicked.connect(self.exit_app)
        self.func_button_ppt = QPushButton("PPT")
        self.func_button_ppt.clicked.connect(self.ppt_add_picture)
        self.func_button_add_row = QPushButton("add row")
        self.func_button_add_row.clicked.connect(self.add_row)
        self.func_button_delete_row = QPushButton("delete row")
        self.func_button_delete_row.clicked.connect(self.delete_row)
        self.func_button_update_row = QPushButton("update row")
        self.func_button_update_row.clicked.connect(self.update_row)
        self.func_button_layout.addWidget(self.func_button_clear,0,0,1,1)
        self.func_button_layout.addWidget(self.func_button_ppt,0,1,1,1)
        self.func_button_layout.addWidget(self.func_button_add_row,0,2,1,1)
        self.func_button_layout.addWidget(self.func_button_delete_row,0,3,1,1)
        self.func_button_layout.addWidget(self.func_button_update_row,1,1,1,1)
        self.func_button_layout.addWidget(self.func_button_quit,1,3,1,1)
        self.func_button_layout.addWidget(self.func_button_update_row,1,3,1,1)
        self.func_button_group.setLayout(self.func_button_layout)
        layout.addWidget(self.func_button_group,11,1,1,1)

        # infomation
        self.info_group = QGroupBox('information')
        self.info_group.setFixedWidth(400)
        self.info_layout = QGridLayout()
        self.info_le = []
        self.info_label = ['project name','Phase','[Doc] Rev.','[Doc] Rev. Date','[Doc] Rev. EO','[Doc] Rev. name',
                        'HW&SW Ver.','OEM P/No','Mando P/No','MHE P/No','NSR No','GP 스티커','working Date']
        for idx, self.lbs in enumerate(self.info_label):
            self.info_layout.addWidget(QLabel(self.info_label[idx],self),(idx),0,1,1)
            self.info_le.append(QLineEdit())
            self.info_le[idx].setFixedHeight(20)
            self.info_layout.addWidget(self.info_le[idx],(idx),1,1,2)
        self.info_group.setLayout(self.info_layout)
        layout.addWidget(self.info_group,4,2,8,1)

        # documentation button
        self.doc_button_group = QGroupBox('documentaion button')
        self.doc_button_group.setFixedWidth(250)
        self.doc_button_layout = QVBoxLayout()
        self.doc_button = []
        self.doc_button_name = ['Inspection\nAgreement','Inspection\nReport','Report','Report','Report',
                                'Report','Report','Report','Report','Report','Report']
        for idx, self.doc_button_name in enumerate(self.doc_button_name):
            self.doc_button.append(QPushButton(self.doc_button_name))
            self.doc_button[idx].setFixedHeight(35)
            self.doc_button_layout.addWidget(self.doc_button[idx])
        self.doc_button[0].clicked.connect(self.doc_agreement_build)
        self.doc_button[1].clicked.connect(self.doc_sample_notice_build)
        self.doc_button_group.setLayout(self.doc_button_layout)
        layout.addWidget(self.doc_button_group,4,3,8,1)

        # data table
        self.data_table_group = QGroupBox('data table')
        #self.data_table_group.setFixedWidth(1100)
        self.data_table_layout = QVBoxLayout()
        self.data_table_pjt = QTableWidget(self)
        self.data_table_pjt.itemClicked.connect(self.item_clicked)
        self.data_table_pjt.itemActivated.connect(self.item_clicked)
        self.data_table_pjt.itemSelectionChanged.connect(self.item_clicked)
        self.data_table_layout.addWidget(self.data_table_pjt)
        self.data_table_group.setLayout(self.data_table_layout)
        layout.addWidget(self.data_table_group,0,1,4,4)
        
        self.df = pd.read_csv('test.csv',encoding='utf-8')
        self.df = self.df.drop("Unnamed: 0",1)

        self.data_table_pjt.setRowCount(len(self.df.index))
        self.data_table_pjt.setColumnCount(len(self.df.columns))
        self.data_table_pjt.setHorizontalHeaderLabels(self.df.columns)
        
        for i in range(len(self.df.index)):
            for j in range(len(self.df.columns)):
                self.data_table_pjt.setItem(i, j, QTableWidgetItem(str(self.df.iloc[i, j])))

        # 데이터프레임의 데이터를 list에 담아둔다.
        self.data = []
        for i in range(list(self.df.shape)[0]):
            temp = list(self.df.iloc[i,:])
            self.data.append(temp)
        self.all_records = self.data
        self.all_columns = list(self.df.columns)

        # GUI 
        self.setGeometry(50, 50, 800, 800)
        self.setWindowTitle("Drawing Analyzer")
        self.setLayout(layout)
        self.show() 
        
    def selectButtonClicked(self):
        self.fname = QFileDialog.getOpenFileName(self, 'Open file', './source/')
        if self.fname:
            self.drawing_loc_le.setText(self.fname[0])
            print(self.fname)

        else:
            QMessageBox.about(self, "Warning", "파일을 선택하지 않았습니다.")
            return

    def changeToImageButtonClicked(self):
        pages = convert_from_path(self.fname[0], dpi=int(self.drawing_combo.currentText()))    #400
        
        png_files = glob.glob('result/*.png')

        for png_file in png_files:
            try:
                os.remove(png_file)
            except:
                print(f"Error:{e.strerror}")
        
        for i, page in enumerate(pages):
            filename = os.path.basename(self.fname[0])[:-4]
            page.save(f'{img_dir+filename}_page{i+1:0>2d}.png','PNG')
            self.terminal_browser.append(f'{filename}_page{i+1:0>2d}.png saved...')
        self.terminal_browser.append('Done !')

        self.onlyfiles = [ f for f in listdir(img_dir) if isfile(join(img_dir,f)) ]
        self.images = [cv2.imread(file) for file in glob.glob(img_dir + "/*.png")]

        filename = "{}.png".format(os.getpid())

        for n in range(0,len(self.onlyfiles)):
            cv2.imwrite(filename, self.images[n])
            pic = QPixmap(filename)
            pic = pic.scaledToWidth(600)
            self.drawing_label[n].setPixmap(QPixmap(pic))
            os.remove(filename)

    def findingButtonClicked(self):
        filename = "{}.png".format(os.getpid())
        self.finding_name = QFileDialog.getOpenFileName(self, 'Open file', './object/')
        
        if self.finding_name:
            self.finding_le.setText(self.finding_name[0])
            self.obj = cv2.imread(self.finding_name[0]) # 찾으려는 이미지
            cv2.imwrite(filename, self.obj)
            pic = QPixmap(filename)
            pic=pic.scaledToWidth(600)
            self.finding_obj_label.setPixmap(QPixmap(pic))        
            os.remove(filename)
        else:
            QMessageBox.about(self, "Warning", "파일을 선택하지 않았습니다.")
            return

    def analyzeButtonClicked(self):
        #이미지 매칭 검색
        filename = "{}.png".format(os.getpid())
        final_match_val = 0
        for n in range(0,len(self.onlyfiles)):
            for i, method_name in enumerate(methods):
                img_draw = self.images[n].copy()
                method = eval(method_name)
                res = cv2.matchTemplate(self.images[n], self.obj, method)
                min_val, max_val, min_loc, max_loc = cv2.minMaxLoc(res)

                if method in [cv2.TM_SQDIFF, cv2.TM_SQDIFF_NORMED]:
                    top_left = min_loc
                    match_val = min_val
                else:
                    top_left = max_loc
                    match_val = max_val

            if final_match_val < match_val:
                final_match_val = match_val
                final_image = self.images[n][top_left[1]:top_left[1]+self.obj.shape[0], top_left[0]:top_left[0]+self.obj.shape[1]]

        # 이미지 처리 - OCR 하기 좋은 이미지로 변환
        final_image = cv2.cvtColor(final_image, cv2.COLOR_BGR2GRAY)        
        #ret, final_image = cv2.threshold(final_image, 127, 255, cv2.THRESH_TOZERO + cv2.THRESH_OTSU)
        #final_image = cv2.adaptiveThreshold(final_image,255,cv2.ADAPTIVE_THRESH_MEAN_C,cv2.THRESH_BINARY,55,4)
        #client = vision.ImageAnnotatorClient() - Google Vision

        text = pytesseract.image_to_string(final_image, config=custom_config) 
        text = text.replace("\n\n","\n")
        text = text.replace(" \n","")
        self.terminal_browser.append(text)

        cv2.imwrite(filename, final_image)
        pic = QPixmap(filename)
        pic.save(f'{find_dir}find_page{n:0>2d}.png','PNG')
        pic_display=pic.scaledToWidth(600)
        self.finding_match_label.setPixmap(pic_display)
        
        os.remove(filename)

        # fill line editor
        self.info_le[0].setText(text[text.find("1 Pro",):text.find("\n",text.find("1 Pro",),)].replace("1 Project Name ","")) 
        self.info_le[6].setText(text[text.find("2 H/W",):text.find("\n",text.find("2 H/W",),)].replace("2 H/W, S/W Ver. ","")) 
        self.info_le[7].setText(text[text.find("3 OEM",):text.find("\n",text.find("3 OEM",),)].replace("3 OEM P/NO ",""))
        self.info_le[8].setText(text[text.find("4 MAN",):text.find("\n",text.find("4 MAN",),)].replace("4 MANDO ",""))
        self.info_le[9].setText(text[text.find("5 Sup",):text.find("\n",text.find("5 Sup",),)].replace("5 Supplier P/No ",""))

        file = open('info.txt','w')
        file.writelines(str(text.encode('utf-8-sig')))
        file.close()

        print(top_left, final_match_val)
        self.terminal_browser.append(f'top:' + str(top_left[0]) + f'   left:' + str (top_left[1]))     
        self.terminal_browser.append(f'match late:' + str(round(final_match_val*100,4)) + "%")     

    def clear_text(self):
        self.tb.clear()

    def exit_app(self):
        QCoreApplication.instance().quit()
        sys.exit(app.exec_())

    def ppt_add_picture(self):
        self.fname = QFileDialog.getOpenFileName(self, 'Open file', './find/')
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(1)
        pic = slide.shapes.add_picture( self.fname[0],left, top)
        prs.save('text.pptx')

    def add_row(self):
        new_pjt_info = [(le.text() for le in self.info_le)]
        dfNew = pd.DataFrame(new_pjt_info, columns = self.df.columns)
        self.df = self.df.append(dfNew, ignore_index=True)
        self.data_table_pjt.setRowCount(len(self.df.index))
        self.data_table_pjt.setColumnCount(len(self.df.columns))

        for i in range(len(self.df.index)):
            for j in range(len(self.df.columns)):
                self.data_table_pjt.setItem(i, j, QTableWidgetItem(str(self.df.iloc[i, j])))
        
        #self.df = pd.DataFrame([[1,1,1,1,1,1,1,1,1,1,1,1,1],[2,2,2,2,2,2,2,2,2,2,2,2,2],[3,3,3,3,3,3,3,3,3,3,3,3,3]],columns=self.info_label)
        self.df.to_csv('test.csv',encoding='utf-8-sig')
        self.data_table_pjt.resizeRowsToContents()
        self.data_table_pjt.resizeColumnsToContents()

        self.terminal_browser.append("New pjt info. added")
        for le in self.info_le:
            le.clear()

        print(self.df)

    def delete_row(self):
        self.terminal_browser.append("Row["+str(self.data_table_pjt.currentRow()+1)+ "] Deleted")
        self.df = self.df.drop(self.data_table_pjt.currentRow())
        self.df = self.df.reset_index(drop=True)
        self.data_table_pjt.setRowCount(len(self.df.index))
        self.data_table_pjt.setColumnCount(len(self.df.columns))
        for i in range(len(self.df.index)):
            for j in range(len(self.df.columns)):
                self.data_table_pjt.setItem(i, j, QTableWidgetItem(str(self.df.iloc[i, j])))
        self.df.to_csv('test.csv',encoding='utf-8-sig')
        self.data_table_pjt.resizeRowsToContents()
        self.data_table_pjt.resizeColumnsToContents()
        print(self.df)

    def item_clicked(self):
        self.data_table_pjt.selectRow(self.data_table_pjt.currentRow())
        for i, le in enumerate(self.info_le):
            le.setText(self.data_table_pjt.item(self.data_table_pjt.currentRow(),i).text())
        print(self.df)

    def update_row(self):
        new_pjt_info = [(le.text() for le in self.info_le)]
        print(len(self.df.columns))
        for i in range(len(self.df.columns)):
            self.df.iloc[self.data_table_pjt.currentRow(), i] = self.info_le[i].text()

        self.data_table_pjt.setRowCount(len(self.df.index))
        self.data_table_pjt.setColumnCount(len(self.df.columns))

        for i in range(len(self.df.index)):
            for j in range(len(self.df.columns)):
                self.data_table_pjt.setItem(i, j, QTableWidgetItem(str(self.df.iloc[i, j])))
        
        self.df.to_csv('test.csv',encoding='utf-8-sig')
        self.data_table_pjt.resizeRowsToContents()
        self.data_table_pjt.resizeColumnsToContents()

        self.terminal_browser.append("New pjt info. updated")
        for le in self.info_le:
            le.clear()

        print(self.df)


    def closeEvent(self, event):
        sys.exit(app.exec_())

    def doc_agreement_build(self):
        QMessageBox.about(self, "Warning", '주의사항: 파일 생성 전, 생성할 차종 선택 필수')
        fname = QFileDialog.getOpenFileName(self, 'Open file', './template/')
        self.terminal_browser.append(str(fname[0]))
        load_wb = load_workbook(fname[0], data_only=True)
        load_ws = load_wb['as']
        ## 검사 협정서 부분 ##
        load_ws['B2'] = self.info_le[8].text() #MANDO P/N
        load_ws['B3'] = self.info_le[0].text() #Project name
        load_ws['B5'] = self.info_le[3].text() #+' (' +data[4]+')' #EO No. & Rev.date
        load_ws['B6'] = self.info_le[12].text() # Working Date
        load_ws['B7'] = self.info_le[5].text() # Rev. History
        load_ws['B8'] = self.info_le[2].text() # Rev. ver.

        ## 검사 기준서 부분 ##\
        load_ws['G2'] = 'Project Name : ' + self.info_le[0].text()
        load_ws['G3'] = 'H/W, S/W Ver. : ' + self.info_le[6].text()
        load_ws['G4'] = 'MANDO P/N : ' + self.info_le[8].text()
        load_ws['G5'] = 'SUPPLIER P/N : ' + self.info_le[9].text()
        load_ws['G6'] = 'OEM P/N : ' + self.info_le[7].text()

        ## cover 시트 부분 ##
        try:
             load_cover = load_wb['COVER']
             load_cover['D15'] = load_cover['D18'].value
             load_cover['D18'] = load_ws['G3'].value + ' (Rev'+load_ws['B8'].value +')'
        except:
             print('검사협정서 아님')

        if fname[0][-1] == 'X' or fname[0][-1] == 'x' :
            file_save = fname[0][:-5]
            file_type = '.xlsx'
            
        elif fname[0][-1] == 'S' or fname[0][-1] == 's':
            file_save = fname[0][:-4]
            file_type = '.xls'
        self.terminal_browser.append(file_save)
        #date_now = datetime.today().strftime(\"%Y%m%d\")
        
        write_wb = file_save + '_' +self.info_le[0].text() + '_rev' + self.info_le[2].text() +'_' + file_type
        
        print(write_wb)
        load_wb.save(write_wb)
        
        self.terminal_browser.append(write_wb + "saved")

    def doc_sample_notice_build(self):
        QMessageBox.about(self, "Warning", '주의사항: 파일 생성 전, 생성할 차종 선택 필수')
        fname = QFileDialog.getOpenFileName(self)
        self.terminal_browser.append(str(fname))

if __name__ == "__main__":
    app = QApplication.instance()
    if app is None:
        app = QApplication(sys.argv)
    window = MyWindow()
    window.show()
    app.exec_()
    sys.exit(app.exec_())